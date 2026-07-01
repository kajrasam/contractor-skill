from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    # 1. Initialize Presentation
    prs = Presentation()

    # Define common styles
    title_slide_layout = prs.slide_layouts[0]
    bullet_slide_layout = prs.slide_layouts[1]
    
    # ----------------------------------------------------
    # Slide 1: Title Slide
    # ----------------------------------------------------
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "People-Skill Web Application"
    subtitle.text = "System Architecture & Design\nเธฃเธฐเธเธเธเธฃเธฐเน€เธกเธดเธเนเธฅเธฐเธเธฑเธ’เธเธฒเธเธตเธ”เธเธงเธฒเธกเธชเธฒเธกเธฒเธฃเธ–เธเธเธฑเธเธเธฒเธ"
    
    # Style title
    title.text_frame.paragraphs[0].font.name = 'Tahoma'
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 102, 204) # Blue

    # ----------------------------------------------------
    # Slide 2: เธ เธฒเธเธฃเธงเธกเธเธญเธเธฃเธฐเธเธ (System Overview)
    # ----------------------------------------------------
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    title.text = "1. เธ เธฒเธเธฃเธงเธกเธเธญเธเธฃเธฐเธเธ (System Overview)"
    
    body = slide.placeholders[1]
    tf = body.text_frame
    
    p = tf.paragraphs[0]
    p.text = "เธฃเธฐเธเธ People-Skill เธเธทเธญ Web Application เธชเธณเธซเธฃเธฑเธเธเธฃเธดเธซเธฒเธฃเธเธฑเธ”เธเธฒเธฃเธเนเธญเธกเธนเธฅเธเธธเธเธฅเธฒเธเธฃ"
    p.font.size = Pt(24)
    p.font.name = 'Tahoma'
    
    p2 = tf.add_paragraph()
    p2.text = "เน€เธเนเธฒเธซเธกเธฒเธขเธซเธฅเธฑเธ:"
    p2.level = 0
    p2.font.bold = True
    
    p3 = tf.add_paragraph()
    p3.text = "เธเธฃเธฐเน€เธกเธดเธเธเธตเธ”เธเธงเธฒเธกเธชเธฒเธกเธฒเธฃเธ– (Competency Assessment)"
    p3.level = 1
    
    p4 = tf.add_paragraph()
    p4.text = "เธเธฑเธ”เธ—เธณเนเธเธเธเธฑเธ’เธเธฒเธเธธเธเธฅเธฒเธเธฃเธฃเธฒเธขเธเธธเธเธเธฅ (Individual Development Plan - IDP)"
    p4.level = 1
    
    p5 = tf.add_paragraph()
    p5.text = "เธฅเธ”เธเนเธญเธเธงเนเธฒเธเธเธญเธเธ—เธฑเธเธฉเธฐ (Skill Gaps) เนเธซเนเธ•เธฃเธเน€เธเนเธฒเธซเธกเธฒเธขเธญเธเธเนเธเธฃ"
    p5.level = 1

    # ----------------------------------------------------
    # Slide 3: เธชเธ–เธฒเธเธฑเธ•เธขเธเธฃเธฃเธกเน€เธ—เธเนเธเนเธฅเธขเธต (Technology Stack)
    # ----------------------------------------------------
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    title.text = "2. เธชเธ–เธฒเธเธฑเธ•เธขเธเธฃเธฃเธกเน€เธ—เธเนเธเนเธฅเธขเธต (Tech Stack)"
    
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.text = "เธญเธญเธเนเธเธเธ”เนเธงเธข Client-Server Architecture เนเธเธ Cloud-Native"
    
    p1 = tf.add_paragraph()
    p1.text = "Frontend (เธซเธเนเธฒเธเนเธฒเธ):"
    p1.level = 0
    p1.font.bold = True
    p11 = tf.add_paragraph()
    p11.text = "HTML5, Vanilla JavaScript, TailwindCSS (เธชเธงเธขเธเธฒเธกเนเธฅเธฐ Responsive), Chart.js"
    p11.level = 1
    
    p2 = tf.add_paragraph()
    p2.text = "Backend (เธซเธฅเธฑเธเธเนเธฒเธ):"
    p2.level = 0
    p2.font.bold = True
    p22 = tf.add_paragraph()
    p22.text = "Python (Flask Framework) เธเธฑเธ”เธเธฒเธฃ API เนเธฅเธฐ Business Logic"
    p22.level = 1
    
    p3 = tf.add_paragraph()
    p3.text = "Database (เธเธฒเธเธเนเธญเธกเธนเธฅ):"
    p3.level = 0
    p3.font.bold = True
    p33 = tf.add_paragraph()
    p33.text = "Supabase (PostgreSQL) - Cloud-based เธเธฅเธญเธ”เธ เธฑเธขเนเธฅเธฐเธฃเธงเธ”เน€เธฃเนเธง"
    p33.level = 1
    
    p4 = tf.add_paragraph()
    p4.text = "Deployment:"
    p4.level = 0
    p4.font.bold = True
    p44 = tf.add_paragraph()
    p44.text = "GitHub (Source Control) เนเธฅเธฐ Render.com (Cloud Hosting & CI/CD)"
    p44.level = 1

    # ----------------------------------------------------
    # Slide 4: เนเธกเธ”เธนเธฅเนเธฅเธฐเธเธตเน€เธเธญเธฃเนเธซเธฅเธฑเธ (Core Features)
    # ----------------------------------------------------
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    title.text = "3. เนเธกเธ”เธนเธฅเนเธฅเธฐเธเธตเน€เธเธญเธฃเนเธซเธฅเธฑเธ (Core Features)"
    
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.text = "เธเธฑเธเธเนเธเธฑเธเธชเธณเธเธฑเธเธเธญเธเธฃเธฐเธเธ"
    
    features = [
        "Authentication & Role Management: เนเธเนเธเธชเธดเธ—เธเธดเน เธเธเธฑเธเธเธฒเธ, เธซเธฑเธงเธซเธเนเธฒเธเธฒเธ, เนเธญเธ”เธกเธดเธ",
        "Dashboard & Analytics: เธชเธฃเธธเธเธ เธฒเธเธฃเธงเธกเธ”เนเธงเธข Radar Chart (Target vs Actual)",
        "Competency Assessment: เธเธฃเธฐเน€เธกเธดเธเธ—เธฑเธเธฉเธฐเนเธฅเธฐเธเธณเธเธงเธ“ Skill Gap เธญเธฑเธ•เนเธเธกเธฑเธ•เธด",
        "Individual Development Plan (IDP): เนเธเธฐเธเธณเธซเธฅเธฑเธเธชเธนเธ•เธฃเธเธฑเธ’เธเธฒเธเธธเธเธฅเธฒเธเธฃเธ•เธฒเธกเธ—เธฑเธเธฉเธฐเธ—เธตเนเธเธฒเธ”",
        "Admin Control: เธเธฑเธ”เธเธฒเธฃเธเนเธญเธกเธนเธฅเธเธเธฑเธเธเธฒเธเนเธฅเธฐ Export เธเนเธญเธกเธนเธฅเน€เธเนเธ Excel"
    ]
    
    for feat in features:
        p = tf.add_paragraph()
        p.text = feat
        p.level = 1
        p.font.size = Pt(20)

    # ----------------------------------------------------
    # Slide 5: เธเนเธญเธ”เธตเนเธฅเธฐเธเธธเธ”เน€เธ”เนเธ (Key Advantages)
    # ----------------------------------------------------
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    title.text = "4. เธเธธเธ”เน€เธ”เนเธเธเธญเธเธฃเธฐเธเธ (Key Advantages)"
    
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.text = "เธ—เธณเนเธกเธฃเธฐเธเธเธเธตเนเธ–เธถเธเธ•เธญเธเนเธเธ—เธขเนเธญเธเธเนเธเธฃเธขเธธเธเนเธซเธกเน"
    
    advs = [
        "Modern & Interactive UI: เธ”เธตเนเธเธเนเธชเธงเธขเธเธฒเธก เนเธเนเธเธฒเธเธเนเธฒเธข เน€เธซเนเธเธ เธฒเธเธฃเธงเธกเนเธ”เนเธ—เธฑเธเธ—เธต",
        "Cloud-Native & Scalable: เนเธฎเธชเธ•เนเธเธ Cloud เธเธขเธฒเธขเธฃเธฐเธเธเธฃเธญเธเธฃเธฑเธเธเธเธเธณเธเธงเธเธกเธฒเธเนเธ”เนเธเนเธฒเธข",
        "Real-time API & Speed: เธเธฃเธฐเธกเธงเธฅเธเธฅเธเนเธฒเธ REST API เธ—เธณเธเธฒเธเนเธซเธฅเธฅเธทเนเธเธฃเธงเธ”เน€เธฃเนเธง",
        "Actionable Insights: เธฃเธฐเธเธเธเธณเธเธงเธ“ Skill Gap เธญเธฑเธ•เนเธเธกเธฑเธ•เธด เนเธฅเธฐเธ•เนเธญเธขเธญเธ”เธเธฑเธ”เธ—เธณ IDP เนเธ”เนเธ—เธฑเธเธ—เธต",
        "Data-Driven Decisions: เธกเธตเธเนเธญเธกเธนเธฅเธเธฃเธฐเธเธญเธเธเธฒเธฃเธ•เธฑเธ”เธชเธดเธเนเธเธชเธณเธซเธฃเธฑเธเธเธนเนเธเธฃเธดเธซเธฒเธฃ"
    ]
    
    for adv in advs:
        p = tf.add_paragraph()
        p.text = adv
        p.level = 1
        p.font.size = Pt(22)
        
    # Save the presentation
    file_path = "People_Skill_System_Design.pptx"
    prs.save(file_path)
    print(f"Presentation saved as {file_path}")

if __name__ == '__main__':
    create_presentation()
